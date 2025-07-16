# -*- coding: utf-8 -*-
"""

@author: gracedavenport
"""

import fitz
import pptx
import openpyxl
import docx
#%%

def doc_to_string(path):
    """
    Reads a .docx file from the specified path and returns its text content as a single string.

    Parameters:
        path (str): The file path to the .docx document.

    Returns:
        string (str): A single string containing all text in the document.
    """
    
    doc = docx.Document(path)
    
    string = "".join([para.text for para in doc.paragraphs])
    
    return string


def ppt_to_string(path):
    """
    Reads a .pptx file from a specified path and returns its text contents as a single string.

    Parameters:
        path (str): The file path to the .pptx presentation.

    Returns:
        string (str): A single string containing all text from all slides in the presentation.
    """
    
    ppt = pptx.Presentation(path)
    
    string = ""
    
    for slide in ppt.slides:
        
        for shape in slide.shapes:
            
            string += getattr(shape, "text", "")
            
    return string


def excel_to_string(path):
    """
    Reads a .xlsx file from a specified path and returns its text conetents as a single string.

    arameters:
        path (str): The path to the .xlsx file.

    Returns:
        string (str): A single string containing all unique values from all cells in the workbook.
    """
    workbook = openpyxl.load_workbook(path)
    
    all_strings = set()
    
    for sheet_name in workbook.sheetnames:
        
        sheet = workbook[sheet_name]
        
        for row in sheet.iter_rows():
            
            for cell in row:
                
                if cell.value is not None and isinstance(cell.value, str):
                    
                    all_strings.add(cell.value)
                    
    return " ".join(list(all_strings))


def pdf_to_string(path):
    """
    Reads a .pdf file from a specified path and returns its text conetents as a single string.

    arameters:
        path (str): The path to the .pdf file.

    Returns:
        string (str): A single string containing all text in the pdf.
    """
    pdf = fitz.open(path)
    
    string = "".join([page.get_text() for page in pdf])
    
    return string


